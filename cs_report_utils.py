"""
Shared utilities for Customer Summary Report
Used by both webapp (Promagent) and cpainc_webapp (CPAinc).
"""
import io
import json
import os
import time
from datetime import datetime

import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt

# ── US state lookup (used to identify international cities) ──────────────────
_US_STATES = {
    # Two-letter abbreviations
    'AL','AK','AZ','AR','CA','CO','CT','DE','FL','GA','HI','ID','IL','IN','IA',
    'KS','KY','LA','ME','MD','MA','MI','MN','MS','MO','MT','NE','NV','NH','NJ',
    'NM','NY','NC','ND','OH','OK','OR','PA','RI','SC','SD','TN','TX','UT','VT',
    'VA','WA','WV','WI','WY','DC',
    # Full names
    'Alabama','Alaska','Arizona','Arkansas','California','Colorado','Connecticut',
    'Delaware','Florida','Georgia','Hawaii','Idaho','Illinois','Indiana','Iowa',
    'Kansas','Kentucky','Louisiana','Maine','Maryland','Massachusetts','Michigan',
    'Minnesota','Mississippi','Missouri','Montana','Nebraska','Nevada',
    'New Hampshire','New Jersey','New Mexico','New York','North Carolina',
    'North Dakota','Ohio','Oklahoma','Oregon','Pennsylvania','Rhode Island',
    'South Carolina','South Dakota','Tennessee','Texas','Utah','Vermont',
    'Virginia','Washington','West Virginia','Wisconsin','Wyoming',
    'District of Columbia',
    # Territories treated as domestic
    'Puerto Rico','Guam','Virgin Islands','American Samoa',
}

# ── Color palette ─────────────────────────────────────────────────────────────

_CHAIN_COLOR_MAP = {
    'hilton':                                  '#003580',
    'marriott bonvoy':                         '#c41230',
    'hyatt hotels':                            '#8B2635',
    'ihg':                                     '#21a038',
    'omni hotels':                             '#2c4a2e',
    'wyndham hotels & resorts':                '#005EB8',
    'choice hotels and radisson americas':     '#FF6600',
    'independent / other':                     '#6c757d',
    'independent/other':                       '#6c757d',
    'bwh hotels':                              '#003087',
    'loews hotels':                            '#8B4513',
    'sonesta':                                 '#8B0000',
    'two roads hospitality':                   '#4a4a4a',
    'posadas':                                 '#B5651D',
}
_DEFAULT_COLORS = ['#1f77b4','#ff7f0e','#2ca02c','#d62728','#9467bd',
                   '#8c564b','#e377c2','#7f7f7f','#bcbd22','#17becf',
                   '#aec7e8','#ffbb78','#98df8a','#ff9896']

def _chain_color(name, idx):
    return _CHAIN_COLOR_MAP.get((name or '').lower()) or _DEFAULT_COLORS[idx % len(_DEFAULT_COLORS)]


# ── Chart generators ──────────────────────────────────────────────────────────

def _make_pie_png(brand_rows):
    labels = [r['chain'] for r in brand_rows]
    values = [r['revenue'] or 0 for r in brand_rows]
    colors = [_chain_color(l, i) for i, l in enumerate(labels)]
    fig, ax = plt.subplots(figsize=(8, 5))
    wedges, _, autotexts = ax.pie(
        values, labels=None, colors=colors,
        autopct=lambda p: f'{p:.1f}%' if p > 3 else '',
        startangle=90, pctdistance=0.82,
        wedgeprops={'width': 0.55}
    )
    for t in autotexts:
        t.set_fontsize(8)
    ax.legend(wedges, [f'{l}  ${v:,.0f}' for l, v in zip(labels, values)],
              loc='center left', bbox_to_anchor=(1, 0.5), fontsize=8, framealpha=0.9)
    ax.set_title('Revenue by Hotel Chain', fontsize=13, fontweight='bold', pad=14)
    plt.tight_layout()
    buf = io.BytesIO()
    fig.savefig(buf, format='png', dpi=150, bbox_inches='tight')
    plt.close(fig)
    buf.seek(0)
    return buf


def _make_city_bar_png(chain_groups, top_n=14):
    city_data = {}   # display_key → {'revenue': float, 'bookings': int}
    intl_rev  = 0.0
    intl_bkgs = 0

    for data in chain_groups.values():
        for c in data['cities']:
            state    = c.get('state') or ''
            is_intl  = bool(state) and state not in _US_STATES
            if is_intl:
                intl_rev  += c['revenue']
                intl_bkgs += c.get('bookings', 0)
            else:
                key = f"{c['city']}, {state}" if state else c['city']
                if key not in city_data:
                    city_data[key] = {'revenue': 0.0, 'bookings': 0}
                city_data[key]['revenue']  += c['revenue']
                city_data[key]['bookings'] += c.get('bookings', 0)

    if intl_rev > 0:
        city_data['International'] = {'revenue': intl_rev, 'bookings': intl_bkgs}

    items  = sorted(city_data.items(), key=lambda x: x[1]['revenue'], reverse=True)[:top_n]
    items  = list(reversed(items))

    # Label format: "City, ST  (N)" — N is total bookings/meetings in that city
    cities   = [f"{k}  ({v['bookings']})" for k, v in items]
    revenues = [v['revenue'] for _, v in items]

    # figsize: width=12 keeps aspect ratio ≤ 0.50 when scaled to 12.5" wide on
    # the slide (slide total 7.5" – 0.75" title – 0.25" margin = 6.5" available).
    fig, ax = plt.subplots(figsize=(12, min(6.0, max(3.0, len(cities) * 0.44))))
    bars = ax.barh(cities, revenues, color='#1a3a5c', height=0.6)
    max_rev = max(revenues) if revenues else 1
    for bar, rev in zip(bars, revenues):
        ax.text(bar.get_width() + max_rev * 0.01, bar.get_y() + bar.get_height() / 2,
                f'${rev:,.0f}', va='center', fontsize=8)
    ax.set_xlabel('Revenue (USD)', fontsize=9)
    ax.set_title('Top Cities by Revenue', fontsize=13, fontweight='bold')
    ax.xaxis.set_major_formatter(plt.FuncFormatter(lambda x, _: f'${x:,.0f}'))
    ax.tick_params(axis='y', labelsize=8)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    plt.tight_layout()
    buf = io.BytesIO()
    fig.savefig(buf, format='png', dpi=150, bbox_inches='tight')
    plt.close(fig)
    buf.seek(0)
    return buf


def _make_map_png(city_map_data):
    """
    Build a static map PNG: continental-US pins on an OSM tile background.
    A small overlay box in the top-left lists off-map destinations grouped
    as Alaska / Hawaii / [territory name] / International.
    Returns a BytesIO PNG buffer, or None on failure.
    """
    try:
        import math
        import urllib.request
        from PIL import Image, ImageDraw, ImageFont

        ZOOM    = 4
        TILE_SZ = 256

        def _lng_to_ftx(lng):
            return (lng + 180.0) / 360.0 * (2 ** ZOOM)

        def _lat_to_fty(lat):
            lr = math.radians(lat)
            return (1.0 - math.log(math.tan(lr) + 1.0 / math.cos(lr)) / math.pi) / 2.0 * (2 ** ZOOM)

        # Tile range for continental US
        LAT_MAX, LAT_MIN =  50.0,  23.0
        LNG_MIN, LNG_MAX = -128.0, -64.0
        tx0, ty0 = int(_lng_to_ftx(LNG_MIN)), int(_lat_to_fty(LAT_MAX))
        tx1, ty1 = int(_lng_to_ftx(LNG_MAX)), int(_lat_to_fty(LAT_MIN))
        img_w = (tx1 - tx0 + 1) * TILE_SZ
        img_h = (ty1 - ty0 + 1) * TILE_SZ
        base  = Image.new('RGB', (img_w, img_h), (200, 210, 230))

        # Download OSM tiles
        for tx in range(tx0, tx1 + 1):
            for ty in range(ty0, ty1 + 1):
                url = f'https://tile.openstreetmap.org/{ZOOM}/{tx}/{ty}.png'
                try:
                    req = urllib.request.Request(url,
                            headers={'User-Agent': 'ConferenceDirect/1.0'})
                    with urllib.request.urlopen(req, timeout=8) as resp:
                        tile = Image.open(io.BytesIO(resp.read())).convert('RGB')
                    base.paste(tile, ((tx - tx0) * TILE_SZ, (ty - ty0) * TILE_SZ))
                except Exception:
                    pass

        def ll2px(lat, lng):
            return (int((_lng_to_ftx(lng) - tx0) * TILE_SZ),
                    int((_lat_to_fty(lat) - ty0) * TILE_SZ))

        # Continental 48 states + DC only (shown as map pins).
        # AK, HI, territories, and foreign go into the overlay box.
        _CONTINENTAL = (
            _US_STATES
            - {'AK','Alaska','HI','Hawaii'}
            - {'Puerto Rico','Guam','Virgin Islands','American Samoa',
               'PR','GU','VI','AS'}
        )

        # Territory normalisation map (abbreviation or full → display name)
        _TERR = {
            'Puerto Rico':'Puerto Rico','PR':'Puerto Rico',
            'Guam':'Guam','GU':'Guam',
            'Virgin Islands':'Virgin Islands','VI':'Virgin Islands',
            'American Samoa':'American Samoa','AS':'American Samoa',
            'Northern Mariana Islands':'N. Mariana Islands',
            'MP':'N. Mariana Islands',
        }

        domestic = [p for p in city_map_data
                    if p.get('state') and p['state'] in _CONTINENTAL]

        # Build overlay entries:
        #   AK/HI/territories → grouped by name (one row each)
        #   Foreign cities     → individual rows, sorted by count desc
        _AK = {'AK', 'Alaska'}
        _HI = {'HI', 'Hawaii'}
        _TERR_ORDER = ['Alaska', 'Hawaii', 'Puerto Rico', 'Guam',
                       'Virgin Islands', 'American Samoa', 'N. Mariana Islands']
        grouped  = {}   # territory label → total count
        foreign  = []   # (label, count) for individual international cities
        for p in city_map_data:
            state = p.get('state') or ''
            if state and state in _CONTINENTAL:
                continue
            if state in _AK:
                grouped['Alaska'] = grouped.get('Alaska', 0) + p['count']
            elif state in _HI:
                grouped['Hawaii'] = grouped.get('Hawaii', 0) + p['count']
            elif state in _TERR:
                lbl = _TERR[state]
                grouped[lbl] = grouped.get(lbl, 0) + p['count']
            else:
                # International — keep individual city entry
                lbl = f"{p['city']}, {state}" if state else p['city']
                foreign.append((lbl, p['count']))

        box_entries = [(lbl, grouped[lbl]) for lbl in _TERR_ORDER if lbl in grouped]
        for lbl, cnt in grouped.items():
            if lbl not in _TERR_ORDER:
                box_entries.append((lbl, cnt))
        foreign.sort(key=lambda x: x[1], reverse=True)
        box_entries.extend(foreign)

        draw = ImageDraw.Draw(base)
        try:
            font    = ImageFont.truetype('/System/Library/Fonts/Helvetica.ttc', 13)
            font_sm = ImageFont.truetype('/System/Library/Fonts/Helvetica.ttc', 11)
            font_bd = ImageFont.truetype('/System/Library/Fonts/Helvetica.ttc', 12)
        except Exception:
            font = font_sm = font_bd = ImageFont.load_default()

        max_count = max((p['count'] for p in city_map_data), default=1)

        # ── Domestic pins with repulsion ────────────────────────────────────
        positions = [ll2px(p['lat'], p['lng']) for p in domestic]
        SPREAD_D  = 30
        fp = list(positions)
        for _ in range(60):
            moved = False
            for a in range(len(fp)):
                for b in range(a + 1, len(fp)):
                    dx, dy = fp[b][0] - fp[a][0], fp[b][1] - fp[a][1]
                    d = math.hypot(dx, dy)
                    if d < SPREAD_D:
                        push = (SPREAD_D - d) / 2 + 0.5
                        nx, ny = (dx / d, dy / d) if d > 0.01 else (1, 0)
                        fp[a] = (fp[a][0] - nx * push, fp[a][1] - ny * push)
                        fp[b] = (fp[b][0] + nx * push, fp[b][1] + ny * push)
                        moved = True
            if not moved:
                break

        for i, pin in enumerate(domestic):
            px, py = int(fp[i][0]), int(fp[i][1])
            r = int(11 + (pin['count'] / max_count) * 8)
            draw.ellipse([px-r, py-r, px+r, py+r], fill='#1a3a5c', outline='white', width=2)
            lbl = str(pin['count'])
            try:
                bb = draw.textbbox((0, 0), lbl, font=font)
                tw, th = bb[2]-bb[0], bb[3]-bb[1]
            except AttributeError:
                tw, th = draw.textsize(lbl, font=font)
            draw.text((px - tw//2, py - th//2), lbl, fill='white', font=font)
            ox, oy = int(positions[i][0]), int(positions[i][1])
            if math.hypot(px - ox, py - oy) > 6:
                draw.line([px, py, ox, oy], fill='#1a3a5c', width=1)
                draw.ellipse([ox-3, oy-3, ox+3, oy+3], fill='#1a3a5c',
                             outline='white', width=1)

        # ── Overlay box (top-left) for off-map destinations ─────────────────
        if box_entries:
            PAD    = 8
            LINE_H = 20
            BOX_W  = 185
            box_h  = PAD * 2 + LINE_H + LINE_H * len(box_entries)

            # Light grey background with dark-blue border (matches web style)
            draw.rectangle([10, 10, 10 + BOX_W, 10 + box_h],
                           fill=(245, 245, 245), outline='#1a3a5c', width=2)

            # Title row
            draw.text((10 + PAD, 10 + PAD), 'Off-Map Destinations',
                      fill='#1a3a5c', font=font_bd)

            for idx, (label, count) in enumerate(box_entries):
                row_y = 10 + PAD + LINE_H * (idx + 1) + 2
                dot_r = 7
                cx = 10 + PAD + dot_r
                draw.ellipse([cx-dot_r, row_y, cx+dot_r, row_y+dot_r*2],
                             fill='#1a3a5c')
                lbl = str(count)
                try:
                    bb = draw.textbbox((0, 0), lbl, font=font_sm)
                    tw, th = bb[2]-bb[0], bb[3]-bb[1]
                except AttributeError:
                    tw, th = draw.textsize(lbl, font=font_sm)
                draw.text((cx - tw//2, row_y + dot_r - th//2), lbl,
                          fill='white', font=font_sm)
                draw.text((cx + dot_r + 6, row_y + dot_r - 6),
                          label[:24], fill='#1a1a1a', font=font_sm)

        buf = io.BytesIO()
        base.save(buf, format='PNG')
        buf.seek(0)
        return buf
    except Exception:
        return None


# ── Query helpers ─────────────────────────────────────────────────────────────

def _query_customer_summary(conn, customer, date_from, date_to,
                             customer_col='"Account Name"',
                             status_col='"Booking Status"',
                             date_col='"Start Date"',
                             contract_col='"Contracted Amount"',
                             revenue_col='"Revenue"',
                             chain_col='"Chain"',
                             brand_col='"Brand"',
                             city_col='"City"',
                             state_col='"State"',
                             cast_contract=True,
                             team_member=None,
                             team_member_col=None):
    # customer may be a string or a list of strings
    customer_list = [customer] if isinstance(customer, str) and customer else (customer or [])

    where_parts = [f'({status_col} IS NULL OR {status_col} NOT LIKE \'%Cancel%\')',
                   f'substr({date_col}, 1, 10) BETWEEN ? AND ?']
    params = [date_from, date_to]
    if customer_list:
        placeholders = ','.join('?' * len(customer_list))
        where_parts.append(f'{customer_col} IN ({placeholders})')
        params.extend(customer_list)
    if team_member and team_member_col:
        where_parts.append(f'{team_member_col} = ?')
        params.append(team_member)
    wc = ' AND '.join(where_parts)

    if cast_contract:
        rev_expr = f'CASE WHEN CAST({contract_col} AS REAL) > 0 THEN CAST({contract_col} AS REAL) ELSE COALESCE({revenue_col}, 0) END'
    else:
        rev_expr = f'CASE WHEN {contract_col} > 0 THEN {contract_col} ELSE COALESCE({revenue_col}, 0) END'

    chain_expr = f'COALESCE(NULLIF({chain_col},\'\'), NULLIF({brand_col},\'\'), \'Unknown\')'

    brand_rows = [dict(r) for r in conn.execute(f'''
        SELECT {chain_expr} as chain, COUNT(*) as bookings, SUM({rev_expr}) as revenue
        FROM ReportPipeline r WHERE {wc} GROUP BY chain ORDER BY revenue DESC
    ''', params).fetchall()]

    detail_rows = conn.execute(f'''
        SELECT {chain_expr} as chain,
               COALESCE(NULLIF({city_col},\'\'), \'Unknown\') as city,
               {state_col} as state,
               COUNT(*) as bookings, SUM({rev_expr}) as revenue
        FROM ReportPipeline r WHERE {wc}
        GROUP BY chain, city ORDER BY chain, revenue DESC
    ''', params).fetchall()

    grand_total    = sum(r['revenue'] or 0 for r in brand_rows)
    grand_bookings = sum(r['bookings'] for r in brand_rows)

    chain_groups = {}
    for r in detail_rows:
        ch = r['chain']
        if ch not in chain_groups:
            chain_groups[ch] = {'bookings': 0, 'revenue': 0.0, 'cities': []}
        chain_groups[ch]['bookings'] += r['bookings']
        chain_groups[ch]['revenue']  += r['revenue'] or 0
        chain_groups[ch]['cities'].append({'city': r['city'], 'state': r['state'] or '',
                                            'bookings': r['bookings'], 'revenue': r['revenue'] or 0})
    chain_groups = dict(sorted(chain_groups.items(), key=lambda x: x[1]['revenue'], reverse=True))
    return brand_rows, chain_groups, grand_total, grand_bookings


def _query_meeting_summary(conn, customer, date_from, date_to,
                            customer_col='"Account Name"',
                            status_col='"Booking Status"',
                            date_col='"Start Date"',
                            contract_col='"Contracted Amount"',
                            revenue_col='"Revenue"',
                            event_col='"Event Name"',
                            booking_col='"Booking Id"',
                            hotel_col='"Customer"',
                            city_col='"City"',
                            state_col='"State"',
                            cast_contract=True,
                            team_member=None,
                            team_member_col=None):
    customer_list = [customer] if isinstance(customer, str) and customer else (customer or [])

    where_parts = [f'({status_col} IS NULL OR {status_col} NOT LIKE \'%Cancel%\')',
                   f'substr({date_col}, 1, 10) BETWEEN ? AND ?']
    params = [date_from, date_to]
    if customer_list:
        placeholders = ','.join('?' * len(customer_list))
        where_parts.append(f'{customer_col} IN ({placeholders})')
        params.extend(customer_list)
    if team_member and team_member_col:
        where_parts.append(f'{team_member_col} = ?')
        params.append(team_member)
    wc = ' AND '.join(where_parts)

    if cast_contract:
        rev_expr = f'CASE WHEN CAST({contract_col} AS REAL) > 0 THEN CAST({contract_col} AS REAL) ELSE COALESCE({revenue_col}, 0) END'
    else:
        rev_expr = f'CASE WHEN {contract_col} > 0 THEN {contract_col} ELSE COALESCE({revenue_col}, 0) END'

    rows = conn.execute(f'''
        SELECT COALESCE(NULLIF({city_col},\'\'), \'Unknown\') as city,
               {state_col} as state,
               {event_col} as event_name,
               {booking_col} as booking_id,
               substr({date_col}, 1, 10) as start_date,
               {hotel_col} as hotel,
               {rev_expr} as revenue
        FROM ReportPipeline r WHERE {wc}
        ORDER BY city, {date_col}
    ''', params).fetchall()

    city_meetings = {}
    for r in rows:
        city = r['city']
        if city not in city_meetings:
            city_meetings[city] = {'state': r['state'] or '', 'total_revenue': 0.0, 'meetings': []}
        city_meetings[city]['total_revenue'] += r['revenue'] or 0
        city_meetings[city]['meetings'].append({
            'event_name': r['event_name'] or '—',
            'booking_id': r['booking_id'],
            'start_date': r['start_date'] or '',
            'hotel':      r['hotel'] or '—',
            'revenue':    r['revenue'] or 0,
        })
    city_meetings = dict(sorted(city_meetings.items(), key=lambda x: x[1]['total_revenue'], reverse=True))
    return city_meetings


# ── City Geocoding (cached) ───────────────────────────────────────────────────

_GEO_CACHE_FILE = os.path.join(os.path.dirname(__file__), '.geocache.json')
_GEO_CACHE: dict = {}
_GEO_CACHE_LOADED = False


def _load_geocache():
    global _GEO_CACHE, _GEO_CACHE_LOADED
    if not _GEO_CACHE_LOADED:
        if os.path.exists(_GEO_CACHE_FILE):
            try:
                with open(_GEO_CACHE_FILE) as f:
                    _GEO_CACHE = json.load(f)
            except Exception:
                _GEO_CACHE = {}
        _GEO_CACHE_LOADED = True


def _save_geocache():
    try:
        with open(_GEO_CACHE_FILE, 'w') as f:
            json.dump(_GEO_CACHE, f, indent=2)
    except Exception:
        pass


def geocode_cities(city_state_list):
    """
    Geocode a list of (city, state) tuples to lat/lng.
    Results are cached in .geocache.json — first call may be slow,
    subsequent calls for the same cities are instant.
    Returns dict: {(city, state): [lat, lng]} — missing entries omitted.
    """
    _load_geocache()
    results = {}
    need = []

    for city, state in city_state_list:
        key = f'{city}|{state or ""}'
        if key in _GEO_CACHE:
            val = _GEO_CACHE[key]
            if val:
                results[(city, state)] = val
        else:
            need.append((city, state, key))

    if need:
        try:
            from geopy.geocoders import Nominatim
            geo = Nominatim(user_agent='conference_direct_report_v1', timeout=10)
            for city, state, key in need:
                query = f'{city}, {state}, USA' if state else f'{city}, USA'
                try:
                    time.sleep(1.1)   # Nominatim hard rate-limit: 1 req/sec
                    loc = geo.geocode(query)
                    val = [loc.latitude, loc.longitude] if loc else None
                except Exception:
                    val = None
                _GEO_CACHE[key] = val
                if val:
                    results[(city, state)] = val
            _save_geocache()
        except ImportError:
            pass   # geopy not installed — map will be skipped

    return results


def build_city_map_data(city_meetings):
    """
    Build list of map pin dicts from city_meetings structure.
    Each dict: {city, state, count, revenue, lat, lng}
    """
    pairs = [(city, data.get('state', '')) for city, data in city_meetings.items()
             if city and city.lower() != 'unknown']
    coords = geocode_cities(pairs)
    pins = []
    for city, data in city_meetings.items():
        state = data.get('state', '')
        coord = coords.get((city, state))
        if coord:
            pins.append({
                'city':    city,
                'state':   state,
                'count':   len(data.get('meetings', [])),
                'revenue': round(data.get('total_revenue', 0), 2),
                'lat':     coord[0],
                'lng':     coord[1],
            })
    return pins


# ── Word document builder ─────────────────────────────────────────────────────

def _build_word_doc(customer, date_from, date_to, brand_rows, chain_groups,
                    grand_total, grand_bookings, city_meetings=None, team_member=None):
    from docx import Document
    from docx.shared import Inches, Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement

    def shade_cell(cell, hex_fill):
        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()
        shd = OxmlElement('w:shd')
        shd.set(qn('w:val'), 'clear')
        shd.set(qn('w:color'), 'auto')
        shd.set(qn('w:fill'), hex_fill)
        tcPr.append(shd)

    doc = Document()
    sec = doc.sections[0]
    sec.left_margin = sec.right_margin = Inches(1)
    sec.top_margin = sec.bottom_margin = Inches(0.75)

    t = doc.add_heading('Customer Summary Report', 0)
    t.alignment = WD_ALIGN_PARAGRAPH.CENTER
    sub = doc.add_paragraph()
    sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
    if isinstance(customer, list):
        cust_display = ', '.join(customer) if customer else 'All Customers'
    else:
        cust_display = customer or 'All Customers'
    subtitle_parts = [f'Customer: {cust_display}', f'Period: {date_from} – {date_to}']
    if team_member:
        subtitle_parts.insert(1, f'Team Member: {team_member}')
    r = sub.add_run('   |   '.join(subtitle_parts))
    r.font.size = Pt(11); r.font.color.rgb = RGBColor(0x44, 0x44, 0x44)
    gp = doc.add_paragraph()
    gp.alignment = WD_ALIGN_PARAGRAPH.CENTER
    rg = gp.add_run(f'Generated: {datetime.today().strftime("%B %d, %Y")}')
    rg.font.size = Pt(9); rg.font.color.rgb = RGBColor(0x88, 0x88, 0x88)
    doc.add_paragraph()

    def add_dark_header_row(tbl, headers):
        hdr_cells = tbl.rows[0].cells
        for i, h in enumerate(headers):
            hdr_cells[i].text = h
            hdr_cells[i].paragraphs[0].runs[0].font.bold = True
            hdr_cells[i].paragraphs[0].runs[0].font.size = Pt(10)
            shade_cell(hdr_cells[i], '1A3A5C')
            hdr_cells[i].paragraphs[0].runs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # Section 1 – Pie chart + brand table
    doc.add_heading('Revenue by Hotel Chain', 1)
    pie_buf = _make_pie_png(brand_rows)
    doc.add_picture(pie_buf, width=Inches(5.5))
    doc.paragraphs[-1].alignment = WD_ALIGN_PARAGRAPH.CENTER
    doc.add_paragraph()

    tbl = doc.add_table(rows=1, cols=4)
    tbl.style = 'Table Grid'
    add_dark_header_row(tbl, ['Hotel Chain', '# Bookings', 'Revenue (USD)', '% of Total'])
    for r in brand_rows:
        row = tbl.add_row().cells
        pct = (r['revenue'] / grand_total * 100) if grand_total else 0
        for i, v in enumerate([r['chain'], str(r['bookings']),
                                f"${r['revenue']:,.2f}", f"{pct:.1f}%"]):
            row[i].text = v
            row[i].paragraphs[0].runs[0].font.size = Pt(9)
    row = tbl.add_row().cells
    for i, v in enumerate(['TOTAL', str(grand_bookings), f'${grand_total:,.2f}', '100.0%']):
        row[i].text = v
        row[i].paragraphs[0].runs[0].font.bold = True
        row[i].paragraphs[0].runs[0].font.size = Pt(9)
        shade_cell(row[i], 'E8E8E8')
    doc.add_paragraph()

    # Section 2 – Bar chart + chain/city table
    doc.add_heading('Revenue by Chain and City', 1)
    bar_buf = _make_city_bar_png(chain_groups)
    doc.add_picture(bar_buf, width=Inches(5.5))
    doc.paragraphs[-1].alignment = WD_ALIGN_PARAGRAPH.CENTER
    doc.add_paragraph()

    tbl2 = doc.add_table(rows=1, cols=4)
    tbl2.style = 'Table Grid'
    add_dark_header_row(tbl2, ['Chain / City', '# Bookings', 'Revenue (USD)', '% of Chain'])
    for chain, data in chain_groups.items():
        row = tbl2.add_row().cells
        for i, v in enumerate([chain, str(data['bookings']), f"${data['revenue']:,.2f}", '']):
            row[i].text = v
            row[i].paragraphs[0].runs[0].font.bold = True
            row[i].paragraphs[0].runs[0].font.size = Pt(9)
            shade_cell(row[i], 'D6E4F0')
        for c in data['cities']:
            row = tbl2.add_row().cells
            city_lbl = f"  {c['city']}, {c['state']}" if c['state'] else f"  {c['city']}"
            pct = (c['revenue'] / data['revenue'] * 100) if data['revenue'] else 0
            for i, v in enumerate([city_lbl, str(c['bookings']),
                                    f"${c['revenue']:,.2f}", f"{pct:.1f}%"]):
                row[i].text = v
                row[i].paragraphs[0].runs[0].font.size = Pt(9)
    row = tbl2.add_row().cells
    for i, v in enumerate(['GRAND TOTAL', str(grand_bookings), f'${grand_total:,.2f}', '']):
        row[i].text = v
        row[i].paragraphs[0].runs[0].font.bold = True
        row[i].paragraphs[0].runs[0].font.size = Pt(9)
        shade_cell(row[i], 'E8E8E8')

    # Section 3 – Meeting summary by city
    if city_meetings:
        doc.add_paragraph()
        doc.add_heading('Meeting Summary by City', 1)
        tbl3 = doc.add_table(rows=1, cols=5)
        tbl3.style = 'Table Grid'
        add_dark_header_row(tbl3, ['City / Meeting', 'Start Date', 'Hotel', 'Booking ID', 'Revenue (USD)'])
        for city, data in city_meetings.items():
            city_lbl = f"{city}, {data['state']}" if data['state'] else city
            row = tbl3.add_row().cells
            for i, v in enumerate([city_lbl, '', '', '', f"${data['total_revenue']:,.2f}"]):
                row[i].text = v
                row[i].paragraphs[0].runs[0].font.bold = True
                row[i].paragraphs[0].runs[0].font.size = Pt(9)
                shade_cell(row[i], 'D6E4F0')
            for m in data['meetings']:
                bid = str(int(float(m['booking_id']))) if m['booking_id'] else '—'
                row = tbl3.add_row().cells
                for i, v in enumerate(['  ' + m['event_name'], m['start_date'],
                                        m['hotel'], bid, f"${m['revenue']:,.2f}"]):
                    row[i].text = v
                    row[i].paragraphs[0].runs[0].font.size = Pt(9)
        row = tbl3.add_row().cells
        for i, v in enumerate(['GRAND TOTAL', '', '', '', f'${grand_total:,.2f}']):
            row[i].text = v
            row[i].paragraphs[0].runs[0].font.bold = True
            row[i].paragraphs[0].runs[0].font.size = Pt(9)
            shade_cell(row[i], 'E8E8E8')

    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)
    return buf


# ── PowerPoint builder ────────────────────────────────────────────────────────

def _build_pptx(customer, date_from, date_to, brand_rows, chain_groups,
                grand_total, grand_bookings, city_meetings=None, team_member=None,
                city_map_data=None):
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    DARK_BLUE = RGBColor(0x1a, 0x3a, 0x5c)
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def add_title_bar(slide, text, font_size=22):
        tb = slide.shapes.add_textbox(Inches(0), Inches(0), Inches(13.33), Inches(0.65))
        tf = tb.text_frame
        tf.text = text
        p = tf.paragraphs[0]
        p.runs[0].font.size = Pt(font_size)
        p.runs[0].font.bold = True
        p.runs[0].font.color.rgb = DARK_BLUE

    # Slide 1: Title
    s1 = prs.slides.add_slide(prs.slide_layouts[0])
    s1.shapes.title.text = 'Customer Summary Report'
    tm_line = f'Team Member: {team_member}\n' if team_member else ''
    if isinstance(customer, list):
        cust_display = ', '.join(customer) if customer else 'All Customers'
    else:
        cust_display = customer or 'All Customers'
    s1.placeholders[1].text = (f'{cust_display}\n'
                                f'{tm_line}'
                                f'{date_from}  –  {date_to}\n'
                                f'Generated: {datetime.today().strftime("%B %d, %Y")}')

    # Slide 2: Pie chart + summary table
    s2 = prs.slides.add_slide(blank)
    add_title_bar(s2, 'Revenue by Hotel Chain')
    pie_buf = _make_pie_png(brand_rows)
    # Omit explicit height — python-pptx auto-computes it from the PNG's natural
    # aspect ratio, ensuring the circle is never stretched into an oval.
    s2.shapes.add_picture(pie_buf, Inches(0.3), Inches(0.75), Inches(6.5))
    n = len(brand_rows) + 2
    tbl = s2.shapes.add_table(n, 3, Inches(7.3), Inches(0.75),
                               Inches(5.7), Inches(min(6.0, n * 0.32 + 0.3))).table
    for ci, h in enumerate(['Hotel Chain', 'Revenue', '%']):
        tbl.cell(0, ci).text = h
        tbl.cell(0, ci).text_frame.paragraphs[0].runs[0].font.bold = True
        tbl.cell(0, ci).text_frame.paragraphs[0].runs[0].font.size = Pt(10)
    for ri, r in enumerate(brand_rows, 1):
        pct = (r['revenue'] / grand_total * 100) if grand_total else 0
        for ci, v in enumerate([r['chain'], f"${r['revenue']:,.0f}", f"{pct:.1f}%"]):
            tbl.cell(ri, ci).text = v
            tbl.cell(ri, ci).text_frame.paragraphs[0].runs[0].font.size = Pt(9)
    ti = len(brand_rows) + 1
    for ci, v in enumerate(['TOTAL', f'${grand_total:,.0f}', '100%']):
        tbl.cell(ti, ci).text = v
        tbl.cell(ti, ci).text_frame.paragraphs[0].runs[0].font.bold = True
        tbl.cell(ti, ci).text_frame.paragraphs[0].runs[0].font.size = Pt(9)

    # Slide 3: City bar chart
    s3 = prs.slides.add_slide(blank)
    add_title_bar(s3, 'Top Cities by Revenue')
    bar_buf = _make_city_bar_png(chain_groups)
    # Only specify width; height auto-computed to preserve aspect ratio
    s3.shapes.add_picture(bar_buf, Inches(0.4), Inches(0.75), Inches(12.5))

    # Slide 4: Meeting Locations Map
    if city_map_data:
        map_buf = _make_map_png(city_map_data)
        if map_buf:
            s4 = prs.slides.add_slide(blank)
            add_title_bar(s4, 'Meeting Locations')
            # Width fills slide; height auto-computed from PNG aspect ratio
            s4.shapes.add_picture(map_buf, Inches(0.3), Inches(0.75), Inches(12.73))

    # Slides 5+: Chain/City detail table (paginated)
    all_rows = [('header', ['Chain / City', '# Bookings', 'Revenue', '% of Chain'], True)]
    for chain, data in chain_groups.items():
        all_rows.append(('chain', [chain, str(data['bookings']), f"${data['revenue']:,.2f}", ''], True))
        for c in data['cities']:
            city_lbl = f"  {c['city']}, {c['state']}" if c['state'] else f"  {c['city']}"
            pct = (c['revenue'] / data['revenue'] * 100) if data['revenue'] else 0
            all_rows.append(('city', [city_lbl, str(c['bookings']),
                                       f"${c['revenue']:,.2f}", f"{pct:.1f}%"], False))
    all_rows.append(('total', ['GRAND TOTAL', str(grand_bookings), f'${grand_total:,.2f}', ''], True))

    def _add_table_slides(prs, blank, all_rows, title_base, cols, chunk_sz=24):
        page = 1
        for start in range(0, len(all_rows), chunk_sz):
            chunk = all_rows[start:start + chunk_sz]
            title_txt = title_base if page == 1 else f'{title_base} (cont.)'
            slide = prs.slides.add_slide(blank)
            add_title_bar(slide, title_txt)
            nr = len(chunk)
            row_h = min(0.27, 6.5 / nr)
            tbl = slide.shapes.add_table(nr, cols, Inches(0.3), Inches(0.75),
                                         Inches(12.7), Inches(nr * row_h)).table
            for ri, (rtype, vals, bold) in enumerate(chunk):
                for ci, v in enumerate(vals):
                    cell = tbl.cell(ri, ci)
                    cell.text = v
                    if cell.text_frame.paragraphs[0].runs:
                        cell.text_frame.paragraphs[0].runs[0].font.size = Pt(9)
                        cell.text_frame.paragraphs[0].runs[0].font.bold = bold
            page += 1

    _add_table_slides(prs, blank, all_rows, 'Revenue Detail by Chain and City', 4)

    # Meeting summary slides
    if city_meetings:
        mtg_rows = [('header', ['City / Meeting', 'Start Date', 'Hotel', 'Booking ID', 'Revenue'], True)]
        for city, data in city_meetings.items():
            city_lbl = f"{city}, {data['state']}" if data['state'] else city
            mtg_rows.append(('city', [city_lbl, '', '', '', f"${data['total_revenue']:,.2f}"], True))
            for m in data['meetings']:
                bid = str(int(float(m['booking_id']))) if m['booking_id'] else '—'
                mtg_rows.append(('meeting', ['  ' + m['event_name'], m['start_date'],
                                              m['hotel'], bid, f"${m['revenue']:,.2f}"], False))
        mtg_rows.append(('total', ['GRAND TOTAL', '', '', '', f'${grand_total:,.2f}'], True))
        _add_table_slides(prs, blank, mtg_rows, 'Meeting Summary by City', 5)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf
